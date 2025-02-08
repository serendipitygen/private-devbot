from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    PythonLoader
)
from langchain.schema import Document
import tempfile
import os
from pptx import Presentation  # 추가
import time

# PPTX 커스텀 로더 추가
class PPTXLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        prs = Presentation(self.file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return [Document(page_content="\n".join(text), metadata={"source": self.file_path})]

def process_file(file_content, filename, file_path):  # file_path 매개변수 추가
    """
    file_content는 메모리에 있는 바이트 데이터입니다.
    하지만 문서 로더는 파일 시스템의 파일 경로가 필요합니다.
    따라서 임시로 파일을 생성해서 처리합니다.
    """
    tmp = None
    try:
        # 1. 메모리의 내용을 임시 파일로 저장
        with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as tmp:
            tmp.write(file_content)
            tmp.flush()
            tmp_path = tmp.name
        
        print(f"[DEBUG] Processing file - Name: {filename}, Path: {tmp_path}")
        
        # 파일 처리
        file_ext = os.path.splitext(filename)[1].lower()
        print(f"[DEBUG] Processing file type: {file_ext}")

        # 파일 유형별 처리
        if file_ext == ".pdf":
            loader = PyPDFLoader(tmp_path)
        elif file_ext == ".docx":
            loader = Docx2txtLoader(tmp_path)
        elif file_ext == ".pptx":
            loader = PPTXLoader(tmp_path)
        elif file_ext in [".txt", ".md"]:
            loader = TextLoader(tmp_path, encoding='utf-8')
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

        # 문서 로드
        docs = loader.load()
        print(f"[DEBUG] Loaded {len(docs)} documents from {filename}")

        # 메타데이터 추가 시 file_path 사용
        for doc in docs:
            doc.metadata.update({
                "source": file_path,  # filename 대신 file_path 사용
                "filename": filename,  # 원본 파일명도 보존
                "file_type": file_ext[1:],
                "timestamp": time.time()
            })

        return docs

    except Exception as e:
        print(f"[ERROR] Failed to process {filename}: {str(e)}")
        raise

    finally:
        # 3. 임시 파일 정리
        if tmp is not None:
            try:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
                    print(f"[DEBUG] Cleaned up temp file: {tmp_path}")
            except Exception as e:
                print(f"[WARNING] Temp file cleanup failed: {str(e)}")
                # 오류는 기록하지만 진행은 계속함

def get_directory_size(directory: str) -> int:
    total = 0
    for dirpath, dirnames, filenames in os.walk(directory):
        for f in filenames:
            fp = os.path.join(dirpath, f)
            if os.path.exists(fp):
                total += os.path.getsize(fp)
    return total
