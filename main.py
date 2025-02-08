import time
import os
import asyncio
import json
import numpy as np

from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import JSONResponse
from fastapi.responses import Response
from pydantic import BaseModel
from typing import List, Optional
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,  # 추가
    CSVLoader,  # 추가
    PythonLoader,
)
from langchain.docstore.document import Document


from pptx import Presentation

import config
import logger_util

from contextlib import asynccontextmanager
from fastapi.middleware.cors import CORSMiddleware
from vector_store import VectorStore

# 벡터 저장소 및 관련 변수 초기화
vector_store = VectorStore()

logger = logger_util.get_logger()

class PPTXLoader:
    def __init__(self, file_path: str):
        self.file_path = file_path
    
    def load(self) -> list:
        prs = Presentation(self.file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return [Document(page_content="\n".join(text_content), metadata={"source": self.file_path})]


@asynccontextmanager
async def lifespan(app: FastAPI):
    try:
        yield  # FastAPI 애플리케이션 실행
    except asyncio.CancelledError:
        # 종료 시 CancelledError를 무시하도록 처리
        pass
    finally:
        try:
            if vector_store:
                vector_store.save_vector_db()  # 저장 호출
                logger.info("[INFO] Saved vector store before shutdown")
        except Exception as e:
            logger.error(f"[ERROR] Shutdown error: {e}")

# FastAPI 앱 초기화
app = FastAPI(lifespan=lifespan)

# CORS 설정
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class SearchRequest(BaseModel):
    query: str
    k: int = 5

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in config.ALLOWED_EXTENSIONS

@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    file_path: str = Form(...),
):
    global vector_store

    try:
        logger.debug(f"**********[DEBUG] Upload request - Path: {file_path}")
        
        file_contents = await file.read()
        content = await vector_store.upload(file_path=file_path, file_name=file.filename,
                            content=file_contents)
        return JSONResponse(content=content)

    except Exception as e:
        logger.exception(f"[ERROR] Upload failed: {str(e)}")
        raise HTTPException(500, detail=f"Upload failed: {str(e)}")


@app.post("/search")
async def search_documents(request: SearchRequest):
    global vector_store

    try:
        print(f"[DEBUG] Searching with query: {request.query}, {request.k}")
        
        result = vector_store.search(request.query, request.k)
    
        # json.dumps를 사용하여 ensure_ascii=False로 한글이 깨지지 않도록 함.
        json_data = json.dumps(result, ensure_ascii=False)
        return Response(
            content=json_data,
            media_type="application/json; charset=utf-8"
        )
    except Exception as e:
        logger.exception(f"[ERROR] Search failed: {str(e)}")
        raise HTTPException(500, detail=str(e))

@app.get("/documents")
async def get_documents():
    global vector_store

    try:
        documents = vector_store.get_documents()
        
        result = {"documents": sorted(documents, key=lambda x: x["file_name"])}
        return result
    except Exception as e:
        logger.exception(str(e))
        raise HTTPException(500, detail=str(e))

class DeleteRequest(BaseModel):
    file_paths: List[str]

@app.delete("/documents")
async def delete_documents(request: DeleteRequest):
    try:
        vector_store.delete_documents(request.file_paths)
        return {"message": "Documents deleted successfully."}
    except Exception as e:
        logger.exception(str(e))
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/documents/all")
async def delete_all_documents():
    try:
        vector_store.delete_all_documents()
        return {"message": "All documents deleted successfully."}
    except Exception as e:
        logger.exception(str(e))
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/status")
async def get_status():
    try:
        document_count = vector_store.get_indexed_file_count()
        index_size = vector_store.get_db_size() / (1024 * 1024)  # MB로 변환
        return {
            "document_count": document_count,
            "index_size_mb": round(index_size, 2),
            "index_path": vector_store.get_vector_db_path()
        }
    except Exception as e:
        logger.error(str(e))
        raise HTTPException(500, detail=str(e))

@app.post("/reset")
async def reset_storage():
    global vector_store

    try:
        vector_store.empty_vector_store()

        return JSONResponse(content={
            "status": "success", 
            "message": "Vector store and indexed files have been reset."
        }, status_code=200) # 상태 코드 200으로 변경
    except Exception as e:
        logger.exception(f"[ERROR] Reset failed: {str(e)}")
        raise HTTPException(500, detail=f"Reset failed: {str(e)}")

def get_loader_class(file_name):
    extension = file_name.split(".")[-1].lower()
    loader_map = {
        "tsv": CSVLoader,
        "csv": CSVLoader,
        "pdf": PyPDFLoader,
        "txt": TextLoader,
        "md": TextLoader,
        "pptx": PPTXLoader,
        "docx": Docx2txtLoader,
        "py": PythonLoader,
        "html": TextLoader,
    }
    if extension not in loader_map:
        raise ValueError(f"Unsupported file extension: {extension}")
    return loader_map.get(extension)


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=config.server_port)